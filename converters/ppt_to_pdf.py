from pptx import Presentation
from fpdf import FPDF
import os

def convert_ppt_to_pdf(ppt_path, output_folder):
    """
    Converts a PowerPoint presentation to a PDF by extracting text content from slides.

    Parameters:
        ppt_path (str): Path to the PowerPoint (.pptx) file.
        output_folder (str): Path to save the generated PDF.

    Returns:
        str: Path to the saved PDF.

    Raises:
        FileNotFoundError: If the PowerPoint file is not found.
        ValueError: If an error occurs while processing the file.
    """

    # Validate input file and ensure output folder
    if not os.path.isfile(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    os.makedirs(output_folder, exist_ok=True)

    # Load the presentation
    try:
        presentation = Presentation(ppt_path)
    except Exception as e:
        raise ValueError(f"Error loading PowerPoint file: {e}")

    # Initialize PDF
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font('Arial', size=12)

    # Extract text from slides and add to PDF
    for slide in presentation.slides:
        slide_content = []

        # Collect all text in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        # Add content to a new PDF page if any text exists
        if slide_content:
            pdf.add_page()
            for text in slide_content:
                pdf.multi_cell(0, 10, txt=text, border=0, align='L')

    # Save the PDF
    pdf_filename = os.path.splitext(os.path.basename(ppt_path))[0] + '.pdf'
    pdf_path = os.path.join(output_folder, pdf_filename)

    try:
        pdf.output(pdf_path)
    except Exception as e:
        raise ValueError(f"Error saving PDF file: {e}")

    return pdf_path
