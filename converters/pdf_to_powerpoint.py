import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def convert_pdf_to_powerpoint(pdf_path, output_folder):
    """
    Converts a PDF file to a PowerPoint presentation with each page as an image.
    """
    powerpoint_filename = os.path.basename(pdf_path).replace('.pdf', '.pptx')
    powerpoint_path = os.path.join(output_folder, powerpoint_filename)

    # Create a PowerPoint presentation
    presentation = Presentation()

    try:
        # Open the PDF file
        pdf_document = fitz.open(pdf_path)

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]

            # Render the page as an image
            pix = page.get_pixmap(dpi=150)  # Increase DPI for better quality
            image_path = os.path.join(output_folder, f"page_{page_num + 1}.png")
            pix.save(image_path)

            # Add a slide with the image
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank slide layout
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

            # Delete the temporary image after adding it to the slide
            os.remove(image_path)

        # Save the PowerPoint presentation
        presentation.save(powerpoint_path)
    
    except Exception as e:
        raise ValueError(f"Error during PDF to PowerPoint conversion: {e}")

    return powerpoint_path
